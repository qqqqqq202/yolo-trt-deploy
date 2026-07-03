#!/usr/bin/env python3
"""Generate group meeting PPT for Intrusion Detector project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
ACCENT    = RGBColor(0x00, 0xBF, 0xA5)    # teal accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xCC)
RED       = RGBColor(0xFF, 0x44, 0x44)
ORANGE    = RGBColor(0xFF, 0xA5, 0x00)
GREEN     = RGBColor(0x00, 0xE6, 0x76)

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    """Fill slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=LIGHT, line_spacing=1.5, font_name="Microsoft YaHei"):
    """Add a text box with multiple paragraphs (one per line string)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return tf


def add_rect(slide, left, top, width, height, fill_color=ACCENT, text="",
             font_size=11, font_color=WHITE):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, left, top, width, height):
    """Add a right arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x55)
    shape.line.fill.background()
    return shape


def add_dot_line(slide, left, top, text, dot_color=ACCENT, font_size=12):
    """Small dot + label."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.15), Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = dot_color
    shape.line.fill.background()
    add_text_box(slide, left + 0.3, top - 0.02, 6, 0.3, text, font_size=font_size, color=LIGHT)


# ============================================================================
# SLIDE 1 — 封面
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_text_box(slide, 1.5, 1.5, 10, 1.2,
             "基于 YOLOv8 + TensorRT 的人员闯入检测系统",
             font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_text_box(slide, 1.5, 2.8, 10, 0.8,
             "从算法到产品化的 C++ 高性能部署实践",
             font_size=24, color=ACCENT, alignment=PP_ALIGN.LEFT)

# Separator line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.8), Inches(3), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

add_text_box(slide, 1.5, 4.2, 6, 0.6,
             "技术栈：YOLOv8n · TensorRT 11 · CUDA 12.8 · OpenCV · C++17",
             font_size=16, color=LIGHT)

add_text_box(slide, 1.5, 5.0, 6, 0.5,
             "硬件平台：NVIDIA GeForce RTX 5060 Laptop GPU",
             font_size=14, color=LIGHT)

add_text_box(slide, 9.0, 6.5, 4, 0.5,
             "2026.07",
             font_size=14, color=LIGHT, alignment=PP_ALIGN.RIGHT)


# ============================================================================
# SLIDE 2 — 项目背景
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 8, 0.8, "项目背景与目标", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Left: 背景
add_text_box(slide, 0.8, 1.6, 5.5, 0.5, "▎背景", font_size=20, color=ACCENT, bold=True)
add_multiline(slide, 0.8, 2.2, 5.5, 3.5, [
    "• 人员闯入检测是安防监控、智慧园区、工业安全的核心需求",
    "• 算法团队通常止步于 Python 脚本/notebook，缺乏",
    "   工程化交付能力",
    "• 产品化落地的关键瓶颈：推理速度、部署便捷性、",
    "   系统稳定性",
], font_size=16, color=LIGHT, line_spacing=1.8)

# Right: 目标
add_text_box(slide, 7.2, 1.6, 5.5, 0.5, "▎目标", font_size=20, color=GREEN, bold=True)
add_multiline(slide, 7.2, 2.2, 5.5, 3.5, [
    "• 构建完整的 C++ TensorRT 推理管线",
    "• 打通 PyTorch → ONNX → TensorRT Engine 全链路",
    "• 实现 ROI 区域闯入判定 + 实时可视化",
    "• 证明推动算法产品化落地的工程能力",
], font_size=16, color=LIGHT, line_spacing=1.8)

# Bottom highlight box
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "核心命题：如何把一纸算法变成一套可交付的 C++ 推理系统？"
p.font.size = Pt(16); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER


# ============================================================================
# SLIDE 3 — 训练 vs 部署
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 10, 0.8, "训练 vs 部署：联系与区别", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Top: 一句话
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.7))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "训练是在实验室造发动机，部署是把发动机装到车上跑起来 — 模型只是发动机，TensorRT 部署才是整车交付"
p.font.size = Pt(14); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

# Left: 联系
add_text_box(slide, 0.8, 2.5, 5.5, 0.4, "▎联系：部署依赖训练产出", font_size=18, color=ACCENT, bold=True)
add_multiline(slide, 0.8, 3.0, 5.5, 1.5, [
    "训练产出 PyTorch 权重 (.pt, 6.2 MB)",
    "   ↓",
    "部署以 .pt 为起点 → 导出 ONNX → 编译 TensorRT Engine",
    "",
    "➜ 训练精度越高，部署检测越准（唯一传递关系）",
], font_size=14, color=LIGHT, line_spacing=1.4)

# Right: 区别 (table)
add_text_box(slide, 7.2, 2.5, 5.5, 0.4, "▎区别：五个维度完全不同", font_size=18, color=ORANGE, bold=True)

diff_lines = [
    ("", "训练", "部署"),
    ("语言/框架", "Python + PyTorch", "C++ + TensorRT"),
    ("目标", "追求精度 (mAP)", "追求速度 (FPS) 和稳定性"),
    ("GPU 使用", "梯度计算，利用率波动", "纯前向推理，持续满载"),
    ("精度策略", "FP32 全精度", "FP16/INT8 量化 (精度损失 < 0.5%)"),
    ("运行环境", "Python + pip 依赖", "一个二进制 + 一个 engine 文件"),
]
for i, (c1, c2, c3) in enumerate(diff_lines):
    y = 3.0 + i * 0.45
    c = WHITE if i == 0 else LIGHT
    fs = 11 if i == 0 else 10
    add_text_box(slide, 7.2, y, 1.7, 0.4, str(c1) if i > 0 else "", font_size=fs, color=c)
    add_text_box(slide, 7.4, y, 2.3, 0.4, str(c2), font_size=fs, color=c, bold=(i == 0))
    add_text_box(slide, 9.8, y, 3.0, 0.4, str(c3), font_size=fs, color=c, bold=(i == 0),
                 font_name="Consolas" if i > 0 else "Microsoft YaHei")

# Bottom left: 格式转换链
add_text_box(slide, 0.8, 4.8, 5.5, 0.4, "▎模型格式转换链路", font_size=16, color=ACCENT, bold=True)

# Pipeline boxes
chain_blocks = [
    ("yolov8n.pt\n6.2 MB\nPyTorch", RGBColor(0x44, 0x62, 0x80)),
    ("yolov8n.onnx\n13 MB\n中间表示", RGBColor(0x3D, 0x6B, 0x8C)),
    ("yolov8n.engine\n14 MB\nGPU 机器码", RGBColor(0x00, 0x88, 0x7A)),
]
cx = 0.8
for label, color in chain_blocks:
    add_rect(slide, cx, 5.4, 2.8, 1.1, fill_color=color, text=label, font_size=11)
    cx += 3.0
add_arrow(slide, 3.65, 5.75, 0.25, 0.35)
add_arrow(slide, 6.85, 5.75, 0.25, 0.35)

add_text_box(slide, 0.8, 6.7, 12, 0.4,
             'ONNX = 模型「通用语言」，TensorRT = 针对 GPU 的「最终编译」，把计算图编译成优化过的 CUDA kernel',
             font_size=11, color=LIGHT)

# Bottom right: 为什么不能直接用 .pt
add_text_box(slide, 10.0, 4.8, 3.0, 0.4, "▎为什么不能直接用 .pt？", font_size=14, color=RED, bold=True)
add_multiline(slide, 10.0, 5.3, 3.0, 1.8, [
    "1. PyTorch 单次推理 50ms+",
    "   TRT 同样模型仅 6ms",
    "2. Python GIL 限制多线程",
    "3. 无法嵌入 C++ 生产系统",
    "4. 无确定性延迟保障",
    "",
    "TensorRT = kernel 自动调优",
    "  + 算子融合 + 内存复用",
], font_size=9, color=LIGHT, line_spacing=1.3)

# ============================================================================
# SLIDE 4 — 部署指标解读
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 10, 0.8, "部署关键指标解读", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# FPS
add_text_box(slide, 0.8, 1.6, 5.8, 0.4, "▎FPS (Frames Per Second)", font_size=20, color=ACCENT, bold=True)
add_multiline(slide, 0.8, 2.1, 5.8, 2.0, [
    "每秒钟能处理多少帧画面。FPS 越高，可并发的视频路数越多。",
    "",
    "我们的实测：104 FPS @ 768×432",
    "→ 输入视频只有 12 FPS，处理速度是播放速度的 9 倍",
    "→ 意味着一台机器可同时跑 8+ 路 12fps 视频不掉帧",
    "",
    "如果 FPS < 视频帧率 → 掉帧，画面跟不上，安防不可接受",
    "影响 FPS：模型大小 · 输入分辨率 · GPU 算力 · 精度策略",
], font_size=13, color=LIGHT, line_spacing=1.3)

# Latency
add_text_box(slide, 7.2, 1.6, 5.5, 0.4, "▎推理延迟 (Latency)", font_size=20, color=GREEN, bold=True)
add_multiline(slide, 7.2, 2.1, 5.5, 2.0, [
    "单帧从输入到输出的时间，单位 ms。",
    "",
    "我们实测：~6 ms / 帧",
    "",
    "延迟 vs 吞吐 (FPS) 的区别：",
    "  • 延迟：一帧要多久 → 关乎实时性",
    "  • 吞吐：一秒能跑多少帧 → 关乎并发能力",
    "  • batch>1 可提高吞吐，但单帧延迟变高",
    "",
    "安防场景延迟更重要 — 人一闯入就立刻报警",
    "不能等攒够一批再判断 → batch=1",
], font_size=13, color=LIGHT, line_spacing=1.3)

# Model size comparison
add_text_box(slide, 0.8, 4.3, 6, 0.4, "▎模型大小的含义", font_size=20, color=ORANGE, bold=True)
add_multiline(slide, 0.8, 4.8, 6, 2.0, [
    "文件大小 ≠ 推理速度！Engine 比 .pt 大，但快 8 倍：",
    "",
    "  yolov8n.pt     6.2 MB    PyTorch 权重 + 优化器状态",
    "  yolov8n.onnx   13  MB    ONNX 计算图 (完整结构描述)",
    "  yolov8n.engine  14  MB    TRT GPU 机器码 + FP16 权重",
    "",
    "Engine 更大是因为 TRT 把算子展开、kernel 代码也嵌入了",
    "大小只影响磁盘占用和加载时间 (我们 < 0.1s)，不影响推理速度",
], font_size=13, color=LIGHT, line_spacing=1.25)

# Why faster
add_text_box(slide, 7.2, 4.3, 5.5, 0.4, "▎为什么部署比训练快？", font_size=20, color=RED, bold=True)
add_multiline(slide, 7.2, 4.8, 5.5, 2.0, [
    "同一个 GPU，训练一帧 30-50ms，部署只要 6ms：",
    "",
    "训练时 (慢)：",
    "  forward + backward + optimizer.step() → 每步更新参数",
    "  需保留中间激活值给反向传播 → 显存占用大",
    "",
    "部署时 (快)：",
    "  只有 forward → 算一遍，不保留中间值",
    "  conv+bn+relu 融合成一个 kernel → 少调 N 次 GPU",
    "  FP16 精度 → 计算单元吞吐翻倍",
], font_size=13, color=LIGHT, line_spacing=1.2)

# Bottom summary
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.5))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "核心结论：部署的价值不在于精度更高，而在于用更低的延迟、更小的开销、更稳定的服务把模型能力交付出去"
p.font.size = Pt(12); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 5 — 技术架构
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 8, 0.8, "系统架构", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Pipeline blocks
blocks = [
    ("视频输入\n(文件/摄像头)", RGBColor(0x44, 0x62, 0x80)),
    ("预处理\nLetterbox+归一化", RGBColor(0x3D, 0x6B, 0x8C)),
    ("TensorRT\nYOLOv8 推理", RGBColor(0x00, 0x88, 0x7A)),
    ("后处理\nNMS+坐标映射", RGBColor(0x3D, 0x6B, 0x8C)),
    ("闯入判定\nROI+中心点", RGBColor(0xC0, 0x5C, 0x00)),
    ("可视化\n框+告警+写视频", RGBColor(0x44, 0x62, 0x80)),
]

x = 0.5
for i, (label, color) in enumerate(blocks):
    add_rect(slide, x, 1.8, 1.85, 1.2, fill_color=color, text=label, font_size=13)
    if i < len(blocks) - 1:
        add_arrow(slide, x + 1.9, 2.15, 0.25, 0.35)
    x += 2.15

# Detail boxes
add_text_box(slide, 0.8, 3.5, 6, 0.4, "▎关键技术细节", font_size=18, color=ACCENT, bold=True)
add_multiline(slide, 0.8, 4.0, 5.8, 3.2, [
    "模型链路",
    "  PyTorch(.pt) → ONNX(.onnx) → TRT Engine(.engine)",
    "  首次构建 5 秒 + 自动缓存，后续 < 0.1 秒加载",
    "",
    "推理优化",
    "  640×640 输入 · FP16 推理 · CUDA Stream 异步",
    "  预处理/推理/后处理全链路 GPU 加速",
    "",
    "闯入逻辑",
    "  人员 bbox 中心点 vs ROI 矩形 · O(1) 判定",
], font_size=13, color=LIGHT, line_spacing=1.2)

add_text_box(slide, 7.2, 3.5, 5.5, 0.4, "▎工程结构", font_size=18, color=ACCENT, bold=True)
add_multiline(slide, 7.2, 4.0, 5.5, 3.2, [
    "Tensor/",
    "├── model/            # ONNX + TRT Engine",
    "├── scripts/          # 模型导出 (Python)",
    "├── src/",
    "│   ├── main.cpp      # 入口 + CLI",
    "│   ├── engine.cpp    # TRT RAII 封装",
    "│   ├── detector.cpp  # 检测器",
    "│   ├── intrusion.cpp # 闯入判定",
    "│   └── visualizer.cpp# 可视化",
    "├── include/           # 头文件",
    "└── CMakeLists.txt    # 构建系统",
], font_size=12, color=LIGHT, line_spacing=1.15)


# ============================================================================
# SLIDE 4 — TensorRT 部署
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 8, 0.8, "TensorRT 部署核心实现", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Left: Engine RAII
add_text_box(slide, 0.8, 1.6, 5.5, 0.5, "▎Engine RAII 封装", font_size=20, color=ACCENT, bold=True)
add_multiline(slide, 0.8, 2.2, 5.5, 2.5, [
    "class Engine {",
    "  IRuntime*      _runtime;   // TRT 运行时",
    "  ICudaEngine*   _engine;    // 推理引擎",
    "  IExecutionContext* _ctx;  // 执行上下文",
    "",
    "  // 构造函数自动判断：缓存加载 或 ONNX 构建",
    "  Engine(onnxPath, enginePath);",
    "  ~Engine() { delete _ctx; delete _engine; ... }",
    "};",
], font_size=11, color=LIGHT, line_spacing=1.1, font_name="Consolas")

# Right: TRT 11 API migration
add_text_box(slide, 7.2, 1.6, 5.5, 0.5, "▎TensorRT 11 新 API 适配", font_size=20, color=ORANGE, bold=True)
add_multiline(slide, 7.2, 2.2, 5.5, 2.5, [
    "旧 (TRT 10)           →  新 (TRT 11)",
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
    "getBindingIndex()      →  getIOTensorName()",
    "enqueueV2(bindings)    →  setTensorAddress()",
    "                            + enqueueV3(stream)",
    "builder->destroy()     →  delete builder",
    "kEXPLICIT_BATCH        →  默认支持 (废弃)",
    "kFP16 flag             →  自动精度选择",
    "platformHasFastFp16()  →  移除 (自动判断)",
], font_size=11, color=LIGHT, line_spacing=1.2, font_name="Consolas")

# Bottom: key metric cards
metrics = [
    ("104 FPS", "平均推理速度", GREEN),
    ("14 MB", "Engine 大小", ACCENT),
    ("< 0.1s", "二次加载耗时", ORANGE),
    ("3.15M", "模型参数量", RGBColor(0x88, 0x88, 0xCC)),
]
for i, (value, label, color) in enumerate(metrics):
    mx = 1.0 + i * 3.05
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(mx), Inches(5.5), Inches(2.7), Inches(1.5))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    card.line.fill.background()
    tf = card.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = value; p.font.size = Pt(36); p.font.color.rgb = color
    p.font.name = "Consolas"; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label; p2.font.size = Pt(12); p2.font.color.rgb = LIGHT
    p2.font.name = "Microsoft YaHei"; p2.alignment = PP_ALIGN.CENTER


# ============================================================================
# SLIDE 5 — 检测 + 闯入管线
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 8, 0.8, "检测管线 & 闯入判定", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Preprocess
add_text_box(slide, 0.8, 1.6, 5.5, 0.4, "▎预处理 Pipeline", font_size=18, color=ACCENT, bold=True)
add_multiline(slide, 0.8, 2.1, 5.5, 2.2, [
    "① Letterbox 缩放：保持宽高比，短边填充到 640×640",
    "② BGR → RGB 色彩空间转换",
    "③ float32 归一化到 [0, 1]",
    "④ HWC → CHW 通道重排，从 OpenCV Mat 到 GPU float*",
    "",
    "关键：scale + padding 参数注入数据末尾，后处理用于坐标还原",
], font_size=13, color=LIGHT, line_spacing=1.6)

# Postprocess
add_text_box(slide, 7.2, 1.6, 5.5, 0.4, "▎后处理 + NMS", font_size=18, color=ACCENT, bold=True)
add_multiline(slide, 7.2, 2.1, 5.5, 2.2, [
    "YOLOv8 输出：[1, 84, 8400]",
    "  84 = 4 (cx,cy,w,h) + 80 (COCO 分类)",
    "  8400 = 80² + 40² + 20² (特征图网格)",
    "",
    "① 遍历 8400 个 anchor → 提取 bbox + 最大分类置信度",
    "② conf > 0.4 的候选 → 按类别 NMS (IoU > 0.5 抑制)",
    "③ 坐标从 letterbox 空间映射回原始图像",
    "④ 过滤 person 类 (COCO class 0)",
], font_size=13, color=LIGHT, line_spacing=1.25)

# Intrusion logic
add_text_box(slide, 0.8, 4.7, 12.5, 0.4, "▎闯入判定逻辑", font_size=18, color=GREEN, bold=True)
add_multiline(slide, 0.8, 5.2, 12, 1.8, [
    "class IntrusionZone {",
    "    cv::Rect2f _rect;   // ROI 矩形区域",
    "",
    "    bool contains(cv::Point2f pt) { return _rect.contains(pt); }",
    "",
    "    vector<Detection> checkIntrusions(detections) {",
    "        // 遍历人员检测，计算 bbox 中心点是否在 ROI 内",
    "        Point2f centre = (box.x + box.w/2,  box.y + box.h/2);",
    "        if (contains(centre)) → INTRUDER!",
    "    }",
    "};",
], font_size=11, color=LIGHT, line_spacing=1.05, font_name="Consolas")


# ============================================================================
# SLIDE 6 — 性能 & 实测
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 8, 0.8, "性能测试 & 实测结果", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Specs
add_text_box(slide, 0.8, 1.6, 6, 0.4, "▎测试环境", font_size=18, color=ACCENT, bold=True)
add_multiline(slide, 0.8, 2.1, 5.5, 2.5, [
    "GPU:    NVIDIA GeForce RTX 5060 Laptop (8 GB)",
    "CPU:    Intel Core Ultra 7 255HX",
    "CUDA:   12.8 · TensorRT: 11.1",
    "OS:     Ubuntu 24.04 (WSL2)",
    "视频:   768×432 @ 12 fps · 596 帧 · 5.3 MB",
], font_size=14, color=LIGHT, line_spacing=1.8)

# Results table-like layout
add_text_box(slide, 7.2, 1.6, 6, 0.4, "▎性能指标", font_size=18, color=GREEN, bold=True)

results = [
    ("指标", "数值", "备注"),
    ("平均 FPS", "104.5", "含全管线 (预处理+推理+后处理+可视化)"),
    ("纯推理耗时", "~6 ms", "单帧 TensorRT enqueueV3"),
    ("Engine 体积", "14 MB", "ONNX 13 MB → TRT 14 MB (含 FP16 优化)"),
    ("首次构建", "~5 秒", "ONNX parse + builder optimization"),
    ("缓存加载", "< 0.1 秒", "直接 deserialize，近乎即时"),
    ("内存占用", "~500 MB", "GPU 显存 (模型 + 中间张量)"),
]

for i, (col1, col2, col3) in enumerate(results):
    y = 2.2 + i * 0.55
    c = WHITE if i == 0 else LIGHT
    add_text_box(slide, 7.2, y, 1.5, 0.5, col1, font_size=13 if i > 0 else 14,
                 color=c, bold=(i == 0))
    add_text_box(slide, 8.8, y, 1.3, 0.5, col2, font_size=14 if i > 0 else 14,
                 color=GREEN if i > 0 else c, bold=(i == 0),
                 font_name="Consolas" if i > 0 else "Microsoft YaHei")
    add_text_box(slide, 10.2, y, 3.0, 0.5, col3, font_size=12,
                 color=LIGHT)

# Intrusion events
add_text_box(slide, 0.8, 5.0, 6, 0.4, "▎真实闯入检测事件 (4次)", font_size=18, color=RED, bold=True)
add_multiline(slide, 0.8, 5.5, 12, 1.8, [
    "事件 1: Frame 28-48   (21 帧) — 行人进入 ROI 中心区域",
    "事件 2: Frame 196-213 (18 帧) — 骑自行车者穿过 ROI",
    "事件 3: Frame 357-375 (19 帧) — 行人再次进入",
    "事件 4: Frame 529-540 (12 帧) — 行人通过 ROI 边缘",
    "",
    "  → YOLOv8n 准确检测所有人员目标，IntrusionZone 实时判定无漏报",
], font_size=13, color=LIGHT, line_spacing=1.3)


# ============================================================================
# SLIDE 7 — 产品化亮点
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 10, 0.8, "产品化工程亮点", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

highlights = [
    ("🔧", "零 Python 运行时依赖", "纯 C++ 编译为单一可执行文件，部署只需二进制 + engine 文件"),
    ("⚡", "Engine 自动缓存", "首次构建自动序列化到磁盘，后续运行秒级加载，CI/CD 友好"),
    ("🖥️", "灵活运行模式", "支持实时弹窗 / 视频保存 / headless 批处理三种模式"),
    ("🎯", "CLI 参数化", "ROI 区域、置信度、NMS 阈值、输入输出路径均可命令行配置"),
    ("📐", "RAII 资源管理", "Engine/Context/CUDA 内存全部 RAII 封装，零泄漏"),
    ("🔌", "CMake 跨平台构建", "find_package 自动适配 TensorRT/CUDA/OpenCV 路径"),
]

for i, (icon, title, desc) in enumerate(highlights):
    y = 1.6 + i * 0.95
    add_text_box(slide, 1.0, y, 3.5, 0.4, f"{icon}  {title}", font_size=18, color=WHITE, bold=True)
    add_text_box(slide, 1.0, y + 0.4, 11, 0.4, desc, font_size=13, color=LIGHT)


# ============================================================================
# SLIDE 8 — 总结 & 展望
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 0.4, 10, 0.8, "总结与展望", font_size=32, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# Summary
add_text_box(slide, 0.8, 1.6, 5.5, 0.5, "▎已完成", font_size=22, color=GREEN, bold=True)
add_multiline(slide, 0.8, 2.2, 5.5, 3.5, [
    "✅ 打通 PyTorch → ONNX → TensorRT Engine 全链路",
    "✅ 纯 C++ 推理管线，零 Python 运行时依赖",
    "✅ 适配最新 TensorRT 11 新 API",
    "✅ ROI 闯入判定 + 可视化 + 视频保存",
    "✅ 104 FPS 实时推理性能 (RTX 5060)",
    "✅ CMake 跨平台构建，模块化设计",
], font_size=15, color=LIGHT, line_spacing=1.8)

# Future
add_text_box(slide, 7.2, 1.6, 5.5, 0.5, "▎后续规划", font_size=22, color=ORANGE, bold=True)
add_multiline(slide, 7.2, 2.2, 5.5, 3.5, [
    "🔲 DeepSORT 多目标跟踪 (跨帧 ID 关联)",
    "🔲 多路 RTSP 视频流并发推理",
    "🔲 报警联动 (MQTT/HTTP 推送 + 截图上传)",
    "🔲 INT8 量化进一步提速",
    "🔲 Jetson Orin 嵌入式平台移植",
    "🔲 Docker 容器化 + REST API 服务化",
], font_size=15, color=LIGHT, line_spacing=1.8)

# Bottom: take-away
box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
box.line.fill.background()
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Takeaway: 算法产品化 = 模型选型 + 推理引擎 + 工程封装 + 业务逻辑，四个环节缺一不可"
p.font.size = Pt(15); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER


# ============================================================================
# Save
# ============================================================================
output_path = "/home/yd/Tensor/Intrusion_Detector_汇报.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
